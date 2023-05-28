import pandas as pd
from pptx import Presentation

df = pd.read_excel('./docs/cemina_registor.xlsx')
prs = Presentation('./docs/badge_format.pptx')

import copy

def duplicate_slide(prs, org_slide) :
    copied_slide = prs.slides.add_slide(org_slide.slide_layout)
    for shape in org_slide.shapes :
        org_element = shape.element
        new_element = copy.deepcopy(org_element)
        copied_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    for value in org_slide.part.rels :
        if 'notesSlide' not in value.reltype :
            copied_slide.part.rels.get_or_add(
                value.reltype,
                value._target
            )
    return copied_slide

from pptx.util import Pt

nametag_count = 0
for person_count in range(len(df)) :
    if person_count % 4 == 0 :
        new_slide = duplicate_slide(prs, prs.slides[0])
        for shape in new_slide.shapes :
            if nametag_count < len(df) and shape.shape_type == 17:
                if  shape.text == '소속' :
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.text = df.iloc[nametag_count, 1]
                    paragraph.font.size = Pt(24)
                    paragraph.font.bold = True
                elif shape.text == '이름' :
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.text = df.iloc[nametag_count, 0]
                    paragraph.font.size = Pt(60)
                    paragraph.font.bold = True
                    nametag_count += 1

prs.save('badge_print.pptx')