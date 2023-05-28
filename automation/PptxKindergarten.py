from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
prs = Presentation('./docs/kindergarten.pptx')

replace_text = [
    ['바른어린이집', 'gentle kids'],
    ['곧은어린이집', 'together house'],
    ['희망어린이집', 'hope house']
]

slide1 = prs.slides[0]
for shape in slide1.shapes :
    # print(shape.shape_type)
    # if shape.shape_type == 1 or shape.shape_type == 17 :
    #     print(shape.text)
    # if shape.shape_type == 1 and '원아모집 (안내)' in shape.text :
    #     shape.text = 'Invitation kindergarten'
    #     paragraph = shape.text_frame.paragraphs[0]
    #     paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    #     paragraph.font.bold = True
    #     paragraph.font.size = Pt(20)
    #     paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if shape.has_table :
        # print(shape.table.cell(row_idx = 0, col_idx = 0).text)
        # print(shape.table.cell(0, 1).text)
        for cell in shape.table.iter_cells() :
            # print(cell.text)
            for re_text in replace_text :
                if cell.text == re_text[0] :
                    cell.text = re_text[1]
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(11)
                    # paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    break

prs.save('kindergarten_update.pptx')
